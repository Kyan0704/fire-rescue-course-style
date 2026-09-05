#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
消防救援课程风格PPT生成器
根据JSON配置生成符合消防救援专业课程风格的PPT
"""
import json
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ========== 风格常量 ==========
COLOR_DEEP_BLUE = RGBColor(0x0A, 0x2A, 0x54)
COLOR_FIRE_ORANGE = RGBColor(0xFF, 0x6B, 0x00)
COLOR_DARK_BG = RGBColor(0x12, 0x12, 0x12)
COLOR_CARD_A = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_CARD_B = RGBColor(0x1C, 0x1C, 0x1E)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
COLOR_MID_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
COLOR_DIM_GRAY = RGBColor(0x9C, 0xA3, 0xAF)

FONT_CN = "Noto Sans SC"
FONT_EN = "Poppins"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """设置幻灯片纯色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COLOR_WHITE, font_name=FONT_CN,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, height=Emu(20000), color=COLOR_WHITE):
    """添加装饰线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image(slide, path, left, top, width, height):
    """添加图片（如果文件存在）"""
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, left, top, width, height)
            return True
        except Exception:
            pass
    return False


def add_circle_icon(slide, left, top, size, color=COLOR_FIRE_ORANGE):
    """添加圆形图标占位"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ========== 页面生成函数 ==========

def create_cover_slide(prs, data):
    """模板1：封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, COLOR_DEEP_BLUE)

    # 背景图
    bg = data.get("background_image", "")
    if bg and os.path.exists(bg):
        add_image(slide, bg, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        # 遮罩
        mask = add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_DEEP_BLUE)
        mask.fill.fore_color.brightness = 0.0

    # 主标题
    add_textbox(slide, Inches(1.11), Inches(2.5), Inches(11.11), Inches(0.83),
                data.get("title", ""), font_size=40, bold=True,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    # 装饰线
    add_line(slide, Inches(5.97), Inches(3.53), Inches(1.39), Emu(25400), COLOR_WHITE)

    # 副标题
    add_textbox(slide, Inches(1.11), Inches(3.75), Inches(11.11), Inches(0.56),
                data.get("subtitle", ""), font_size=20, bold=False,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    # 演讲者
    add_textbox(slide, Inches(1.11), Inches(6.39), Inches(11.11), Inches(0.42),
                data.get("speaker", ""), font_size=14, bold=False,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    return slide


def create_contents_slide(prs, data):
    """模板2：目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DEEP_BLUE)

    # CONTENTS 标题
    add_textbox(slide, Inches(0.83), Inches(1.11), Inches(11.67), Inches(0.69),
                "CONTENTS", font_size=36, bold=True, color=COLOR_WHITE,
                font_name=FONT_EN, alignment=PP_ALIGN.LEFT)

    items = data.get("items", [])
    start_top = 2.22
    item_height = 0.83
    gap = 0.14

    for i, item in enumerate(items[:5]):
        top = start_top + i * (item_height + gap)
        # 背景条（半透明白色效果用深色替代）
        bar = add_rect(slide, Inches(0.83), Inches(top), Inches(11.67), Inches(item_height),
                       RGBColor(0x1A, 0x3A, 0x64))
        # 图标占位
        add_circle_icon(slide, Inches(1.11), Inches(top + 0.2), Inches(0.44), COLOR_FIRE_ORANGE)
        # 文字
        add_textbox(slide, Inches(1.74), Inches(top + 0.14), Inches(10.42), Inches(0.56),
                    item, font_size=18, bold=True, color=COLOR_WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)

    return slide


def create_chapter_slide(prs, data):
    """模板3：章节过渡页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DEEP_BLUE)

    bg = data.get("background_image", "")
    if bg and os.path.exists(bg):
        add_image(slide, bg, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_DEEP_BLUE)

    # CHAPTER标签
    add_textbox(slide, Inches(0.83), Inches(2.22), Inches(8.33), Inches(0.56),
                data.get("label", ""), font_size=20, bold=True,
                color=COLOR_FIRE_ORANGE)

    # 大标题
    add_textbox(slide, Inches(0.83), Inches(2.78), Inches(9.72), Inches(1.11),
                data.get("title", ""), font_size=36, bold=True,
                color=COLOR_WHITE)

    # 装饰线
    add_line(slide, Inches(0.83), Inches(4.02), Inches(1.11), Emu(25400), COLOR_WHITE)

    # 引言
    add_textbox(slide, Inches(0.83), Inches(4.31), Inches(8.33), Inches(0.56),
                data.get("quote", ""), font_size=18, bold=False,
                color=COLOR_WHITE)

    return slide


def create_image_text_slide(prs, data):
    """模板4：左图右文页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DARK_BG)

    # 页面标题
    add_textbox(slide, Inches(0.83), Inches(1.11), Inches(11.67), Inches(0.7),
                data.get("page_title", ""), font_size=20, bold=True,
                color=COLOR_WHITE)

    # 左侧图片
    img = data.get("image", "")
    if not add_image(slide, img, Inches(0.83), Inches(2.2), Inches(5.6), Inches(4.2)):
        add_rect(slide, Inches(0.83), Inches(2.2), Inches(5.6), Inches(4.2), COLOR_CARD_A)

    # 右侧卡片
    add_rounded_rect(slide, Inches(6.8), Inches(2.2), Inches(5.7), Inches(1.7), COLOR_DEEP_BLUE)
    add_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.1), Inches(0.4),
                data.get("card_title", ""), font_size=18, bold=True, color=COLOR_WHITE)
    add_textbox(slide, Inches(7.1), Inches(2.8), Inches(5.1), Inches(1.0),
                data.get("card_body", ""), font_size=14, bold=False,
                color=COLOR_LIGHT_GRAY)

    # 要点
    points = data.get("points", [])
    if points:
        add_textbox(slide, Inches(6.8), Inches(4.1), Inches(5.7), Inches(0.4),
                    data.get("points_title", "核心特征"), font_size=16, bold=True,
                    color=COLOR_WHITE)

        positions = [
            (6.8, 4.7, 7.2, 4.6, 2.5),
            (9.7, 4.7, 10.1, 4.6, 2.4),
            (6.8, 5.4, 7.2, 5.4, 2.5),
            (9.7, 5.4, 10.1, 5.4, 2.4),
        ]
        for i, point in enumerate(points[:4]):
            if i >= len(positions):
                break
            ix, iy, tx, ty, tw = positions[i]
            add_circle_icon(slide, Inches(ix), Inches(iy), Inches(0.3), COLOR_FIRE_ORANGE)
            title = point.get("title", "")
            desc = point.get("desc", "")
            text = f"{title}\n{desc}" if desc else title
            add_textbox(slide, Inches(tx), Inches(ty), Inches(tw), Inches(0.7),
                        text, font_size=12, bold=False, color=COLOR_LIGHT_GRAY)

    return slide


def create_grid_2x2_slide(prs, data):
    """模板5：四宫格卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DARK_BG)

    add_textbox(slide, Inches(0.83), Inches(1.11), Inches(11.67), Inches(0.7),
                data.get("page_title", ""), font_size=20, bold=True, color=COLOR_WHITE)

    cards = data.get("cards", [])
    positions = [
        (0.83, 2.2), (6.9, 2.2),
        (0.83, 4.7), (6.9, 4.7),
    ]

    for i, card in enumerate(cards[:4]):
        if i >= len(positions):
            break
        left, top = positions[i]
        # 卡片
        add_rounded_rect(slide, Inches(left), Inches(top), Inches(5.6), Inches(2.2), COLOR_DEEP_BLUE)
        # 图标
        icon = card.get("icon", "")
        if not add_image(slide, icon, Inches(left + 0.27), Inches(top + 0.3), Inches(0.7), Inches(0.7)):
            add_circle_icon(slide, Inches(left + 0.27), Inches(top + 0.3), Inches(0.7), COLOR_FIRE_ORANGE)
        # 标题
        add_textbox(slide, Inches(left + 1.07), Inches(top + 0.3), Inches(4.2), Inches(0.4),
                    card.get("title", ""), font_size=20, bold=True, color=COLOR_FIRE_ORANGE)
        # 正文
        add_textbox(slide, Inches(left + 0.27), Inches(top + 0.9), Inches(5.1), Inches(1.1),
                    card.get("body", ""), font_size=14, bold=False, color=COLOR_LIGHT_GRAY)

    return slide


def create_timeline_4col_slide(prs, data):
    """模板6：四列时间线页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DARK_BG)

    add_textbox(slide, Inches(0.83), Inches(1.11), Inches(11.67), Inches(0.7),
                data.get("page_title", ""), font_size=20, bold=True, color=COLOR_WHITE)

    columns = data.get("columns", [])
    col_lefts = [0.83, 3.8, 6.8, 9.7]

    for i, col in enumerate(columns[:4]):
        if i >= len(col_lefts):
            break
        left = col_lefts[i]
        # 卡片
        add_rect(slide, Inches(left), Inches(2.2), Inches(2.8), Inches(4.2), COLOR_CARD_A)
        # 顶部装饰条
        add_rect(slide, Inches(left), Inches(2.2), Inches(2.8), Inches(0.1), COLOR_DEEP_BLUE)
        # 图标
        add_circle_icon(slide, Inches(left + 0.27), Inches(2.5), Inches(0.4), COLOR_FIRE_ORANGE)
        # 标题
        add_textbox(slide, Inches(left + 0.27), Inches(3.1), Inches(2.2), Inches(0.4),
                    col.get("title", ""), font_size=16, bold=True, color=COLOR_WHITE)
        # 时间
        add_textbox(slide, Inches(left + 0.27), Inches(3.5), Inches(2.2), Inches(0.3),
                    col.get("time", ""), font_size=12, bold=False, color=COLOR_DIM_GRAY)
        # 分隔线
        add_line(slide, Inches(left + 0.27), Inches(3.9), Inches(2.2), Emu(12700), COLOR_MID_GRAY)
        # 正文
        add_textbox(slide, Inches(left + 0.27), Inches(4.0), Inches(2.2), Inches(2.1),
                    col.get("body", ""), font_size=12, bold=False, color=COLOR_LIGHT_GRAY)

    return slide


def create_two_column_slide(prs, data):
    """模板7：双栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DARK_BG)

    add_textbox(slide, Inches(0.83), Inches(0.6), Inches(11.67), Inches(0.7),
                data.get("page_title", ""), font_size=20, bold=True, color=COLOR_WHITE)

    for col_idx, col_key in enumerate(["left", "right"]):
        col = data.get(col_key, {})
        left = 0.83 if col_idx == 0 else 6.9
        # 卡片
        add_rect(slide, Inches(left), Inches(1.5), Inches(5.6), Inches(5.3), COLOR_CARD_B)
        # 顶部图片
        img = col.get("image", "")
        if not add_image(slide, img, Inches(left), Inches(1.5), Inches(5.6), Inches(2.5)):
            add_rect(slide, Inches(left), Inches(1.5), Inches(5.6), Inches(2.5), COLOR_DEEP_BLUE)
        # 标题
        add_textbox(slide, Inches(left + 0.17), Inches(4.2), Inches(5.2), Inches(0.4),
                    col.get("title", ""), font_size=18, bold=True, color=COLOR_WHITE)
        # 分隔线
        add_line(slide, Inches(left + 0.17), Inches(4.7), Inches(5.2), Emu(12700), COLOR_MID_GRAY)
        # 要点
        points = col.get("points", [])
        for j, point in enumerate(points[:3]):
            py = 4.9 + j * 0.7
            add_circle_icon(slide, Inches(left + 0.17), Inches(py), Inches(0.3), COLOR_FIRE_ORANGE)
            add_textbox(slide, Inches(left + 0.57), Inches(py - 0.1), Inches(4.9), Inches(0.6),
                        point.get("desc", ""), font_size=14, bold=False, color=COLOR_LIGHT_GRAY)

    return slide


def create_summary_slide(prs, data):
    """模板8A：总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DEEP_BLUE)

    # 大标题
    add_textbox(slide, Inches(1.11), Inches(0.83), Inches(11.11), Inches(0.83),
                data.get("title", ""), font_size=48, bold=True,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    # 内容区
    points = data.get("points", [])
    for i, point in enumerate(points[:3]):
        py = 2.5 + i * 1.11
        add_circle_icon(slide, Inches(1.67), Inches(py), Inches(0.56), COLOR_WHITE)
        add_textbox(slide, Inches(2.36), Inches(py), Inches(9.44), Inches(0.42),
                    point.get("title", ""), font_size=20, bold=True, color=COLOR_FIRE_ORANGE)
        add_textbox(slide, Inches(2.36), Inches(py + 0.42), Inches(9.44), Inches(0.56),
                    point.get("body", ""), font_size=16, bold=False, color=COLOR_LIGHT_GRAY)

    # 致谢
    add_textbox(slide, Inches(0), Inches(5.83), Inches(13.33), Inches(0.56),
                data.get("thanks", "感谢聆听！"), font_size=24, bold=True,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    return slide


def create_end_slide(prs, data):
    """模板8B：结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_DEEP_BLUE)

    bg = data.get("background_image", "")
    if bg and os.path.exists(bg):
        add_image(slide, bg, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLOR_DEEP_BLUE)

    add_textbox(slide, Inches(1.11), Inches(2.5), Inches(11.11), Inches(0.83),
                data.get("course_name", ""), font_size=36, bold=True,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.11), Inches(3.4), Inches(11.11), Inches(0.56),
                data.get("speaker", ""), font_size=20, bold=False,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    add_line(slide, Inches(5.97), Inches(4.16), Inches(1.39), Emu(25400), COLOR_WHITE)

    add_textbox(slide, Inches(1.11), Inches(4.44), Inches(11.11), Inches(0.56),
                data.get("slogan", ""), font_size=16, bold=False,
                color=COLOR_FIRE_ORANGE, alignment=PP_ALIGN.CENTER)

    return slide


# ========== 页面类型映射 ==========
SLIDE_CREATORS = {
    "cover": create_cover_slide,
    "contents": create_contents_slide,
    "chapter": create_chapter_slide,
    "image_text": create_image_text_slide,
    "grid_2x2": create_grid_2x2_slide,
    "timeline_4col": create_timeline_4col_slide,
    "two_column": create_two_column_slide,
    "summary": create_summary_slide,
    "end": create_end_slide,
}


def generate_ppt(input_path, output_path):
    """主生成函数"""
    with open(input_path, "r", encoding="utf-8") as f:
        config = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides_config = config.get("slides", [])
    created = 0

    for slide_cfg in slides_config:
        slide_type = slide_cfg.get("type", "")
        data = slide_cfg.get("data", {})
        creator = SLIDE_CREATORS.get(slide_type)
        if creator:
            creator(prs, data)
            created += 1
        else:
            print(f"警告：未知页面类型 '{slide_type}'，已跳过")

    prs.save(output_path)
    print(f"成功生成PPT：{output_path}")
    print(f"共生成 {created} 页幻灯片")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="消防救援课程风格PPT生成器")
    parser.add_argument("--input", "-i", required=True, help="输入JSON配置文件路径")
    parser.add_argument("--output", "-o", required=True, help="输出PPTX文件路径")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"错误：输入文件不存在：{args.input}")
        return

    generate_ppt(args.input, args.output)


if __name__ == "__main__":
    main()
